from pptx import Presentation
from pptx.util import Inches, Pt
from typing import List
import io

def create_pptx(title: str, slides: List[dict]) -> bytes:
    """
    Create a PowerPoint presentation
    slides: [{'title': 'Slide 1', 'content': 'Content here'}, ...]
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    
    # Content slides
    bullet_slide_layout = prs.slide_layouts[1]
    for slide_data in slides:
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = slide_data['title']
        
        content = slide.shapes.placeholders[1].text_frame
        content.text = slide_data['content']
    
    # Save to bytes
    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    return file_stream.read()
