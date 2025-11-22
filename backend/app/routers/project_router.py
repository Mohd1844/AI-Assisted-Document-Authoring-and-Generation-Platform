from fastapi import APIRouter, Depends, HTTPException, Body
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from pydantic import BaseModel
from typing import List, Optional
from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import io
import re
from ..database import get_db
from ..models import Project, Section, User
from ..auth import get_current_user

router = APIRouter(prefix="/projects", tags=["projects"])

# --- SCHEMAS ---
class SectionCreate(BaseModel):
    title: str
    order_index: int

class SlideCreate(BaseModel):
    title: str
    order_index: int

class ProjectConfig(BaseModel):
    sections: Optional[List[SectionCreate]] = None
    slides: Optional[List[SlideCreate]] = None

class ProjectCreate(BaseModel):
    title: str
    document_type: str  # 'docx' or 'pptx'
    topic: str
    config: Optional[ProjectConfig] = None

class SectionResponse(BaseModel):
    id: int
    title: str
    order_index: int
    class Config:
        from_attributes = True

class ProjectResponse(BaseModel):
    id: int
    title: str
    document_type: str
    topic: str
    sections: List[SectionResponse] = []
    class Config:
        from_attributes = True

# --- Helper functions ---
def clean_markdown_headers(text: str) -> str:
    text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
    return text.replace('**', '').replace('*', '').strip()

def parse_content_to_paragraphs(content: str) -> list:
    return [p.strip() for p in content.split('\n\n') if p.strip()]

# --- CREATE PROJECT WITH SECTIONS ---
@router.post("/", response_model=ProjectResponse)
def create_project(
    project_data: ProjectCreate,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    new_project = Project(
        user_id=current_user.id,
        title=project_data.title,
        document_type=project_data.document_type,
        topic=project_data.topic
    )
    db.add(new_project)
    db.commit()
    db.refresh(new_project)

    # Add sections/slides if present
    if project_data.config:
        if project_data.document_type == "docx" and project_data.config.sections:
            for section in project_data.config.sections:
                new_section = Section(
                    project_id=new_project.id,
                    title=section.title,
                    order_index=section.order_index
                )
                db.add(new_section)
        elif project_data.document_type == "pptx" and project_data.config.slides:
            for slide in project_data.config.slides:
                new_section = Section(
                    project_id=new_project.id,
                    title=slide.title,
                    order_index=slide.order_index
                )
                db.add(new_section)
        db.commit()

    sections_out = db.query(Section).filter(Section.project_id == new_project.id).order_by(Section.order_index).all()
    sections_response = [SectionResponse.model_validate(section) for section in sections_out]
    return ProjectResponse(
        id=new_project.id,
        title=new_project.title,
        document_type=new_project.document_type,
        topic=new_project.topic,
        sections=sections_response
    )

# --- GET PROJECT WITH SECTIONS ---
@router.get("/{project_id}", response_model=ProjectResponse)
def get_project(
    project_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    project = db.query(Project).filter(Project.id == project_id, Project.user_id == current_user.id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    sections = db.query(Section).filter(Section.project_id == project.id).order_by(Section.order_index).all()
    sections_response = [SectionResponse.model_validate(section) for section in sections]
    return ProjectResponse(
        id=project.id,
        title=project.title,
        document_type=project.document_type,
        topic=project.topic,
        sections=sections_response
    )

# --- ADD SECTION ENDPOINT ---
@router.post("/{project_id}/sections", response_model=SectionResponse)
def add_section(
    project_id: int,
    section_data: SectionCreate,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == current_user.id
    ).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    new_section = Section(
        project_id=project_id,
        title=section_data.title,
        order_index=section_data.order_index
    )
    db.add(new_section)
    db.commit()
    db.refresh(new_section)
    return SectionResponse.model_validate(new_section)

# --- FEEDBACK ENDPOINT ---
@router.post("/{project_id}/feedback")
def add_feedback(
    project_id: int,
    section_id: int = Body(...),
    feedback_type: str = Body(...),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    section = db.query(Section).filter(
        Section.id == section_id, Section.project_id == project_id
    ).first()
    if not section:
        raise HTTPException(status_code=404, detail="Section not found")
    return {"message": "Feedback saved"}

# --- COMMENT ENDPOINT ---
@router.post("/{project_id}/comment")
def add_comment(
    project_id: int,
    section_id: int = Body(...),
    comment: str = Body(...),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    section = db.query(Section).filter(
        Section.id == section_id, Section.project_id == project_id
    ).first()
    if not section:
        raise HTTPException(status_code=404, detail="Section not found")
    return {"message": "Comment saved"}

# --- EXPORT DOCX OR PPTX (WITH LINE-BASED BULLET SPLITTING) ---
@router.get("/{project_id}/export")
def export_project(
    project_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    project = db.query(Project).filter(Project.id == project_id, Project.user_id == current_user.id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    sections = db.query(Section).filter(Section.project_id == project.id).order_by(Section.order_index).all()

    if project.document_type == 'docx':
        doc = Document()
        title_head = doc.add_heading(project.title, 0)
        title_head.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        topic_para = doc.add_paragraph(f"Topic: {project.topic}")
        topic_para.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        topic_para.runs[0].font.italic = True
        doc.add_paragraph()

        for section in sections:
            heading = doc.add_heading(section.title, level=1)
            heading.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
            for run in heading.runs:
                run.font.color.rgb = RGBColor(31, 56, 100)
            cleaned_content = clean_markdown_headers(section.content or "[No content]")
            paragraphs = parse_content_to_paragraphs(cleaned_content)
            for para_text in paragraphs:
                paragraph = doc.add_paragraph(para_text)
                paragraph.paragraph_format.line_spacing = 1.3
                paragraph.paragraph_format.space_before = Pt(8)
                paragraph.paragraph_format.space_after = Pt(8)
                paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(12)
                    run.font.name = "Calibri"
            doc.add_paragraph()

        bio = io.BytesIO()
        doc.save(bio)
        bio.seek(0)
        return StreamingResponse(
            iter([bio.getvalue()]),
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f"attachment; filename={project.title or 'document'}.docx"}
        )

    elif project.document_type == 'pptx':
        prs = Presentation()
        for section in sections:
            cleaned_content = clean_markdown_headers(section.content or "")
            # Split by lines for better bullets
            lines = [l for l in cleaned_content.split('\n') if l.strip()]
            slides = []
            current_title = section.title
            current_body = []
            for line in lines:
                if (line.strip().startswith("**") and line.strip().endswith("**")) or line.strip().endswith("?"):
                    if current_body:
                        slides.append((current_title, current_body))
                        current_body = []
                    current_title = line.replace("**", "").strip()
                else:
                    current_body.append(line.strip())
            if current_body:
                slides.append((current_title, current_body))

            # Max lines (bullets) per slide
            max_lines_per_slide = 5
            for heading, body_lines in slides:
                chunks = [body_lines[i:i+max_lines_per_slide] for i in range(0, len(body_lines), max_lines_per_slide)]
                for idx, chunk in enumerate(chunks):
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.title.text = heading if idx == 0 else f"{heading} (cont'd)"
                    content_shape = slide.placeholders[1]
                    text_frame = content_shape.text_frame
                    text_frame.word_wrap = True
                    text_frame.clear()
                    # Each line is a bullet:
                    for bullet in chunk:
                        p = text_frame.add_paragraph()
                        p.text = bullet
                        p.font.size = Pt(16)
                        p.font.name = "Calibri"
                        p.alignment = PP_ALIGN.LEFT
                        p.level = 0  # Bullet

        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        return StreamingResponse(
            iter([pptx_bytes.getvalue()]),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={project.title or 'presentation'}.pptx"}
        )
    else:
        raise HTTPException(status_code=400, detail="Unsupported document type")
